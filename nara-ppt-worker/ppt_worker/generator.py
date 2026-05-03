"""나라장터 제안서 PPT 생성기.

슬라이드 디자인 규칙:
- 기본 폰트: 나눔고딕 (Bold/Regular)
- 제목: 24pt Bold, 본문: 12pt Regular
- 메인 컬러: #2B579A, 포인트: #217346
- 슬라이드: 와이드스크린 16:9
"""

from __future__ import annotations

import io
import re
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Literal

import httpx
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

from .theme import Theme, DEFAULT_THEME, get_theme, resolve_theme

# ──────────────────────────── 디자인 상수 ────────────────────────────
# 색/폰트는 Theme로 이전됨. 아래 상수는 레이아웃 수치와 폴백용으로만 유지.

TITLE_SIZE = Pt(24)
BODY_SIZE = Pt(12)
CAPTION_SIZE = Pt(10)
PAGE_NUM_SIZE = Pt(9)

COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_WIDTH = Emu(12192000)   # 16:9 와이드스크린
SLIDE_HEIGHT = Emu(6858000)

# 공통 마진
MARGIN_LEFT = Inches(0.7)
MARGIN_RIGHT = Inches(0.7)
CONTENT_WIDTH = Inches(8.9)   # 10.33 - 0.7*2
HEADER_TOP = Inches(0.4)
BODY_TOP = Inches(1.5)
BODY_HEIGHT = Inches(5.0)
FOOTER_TOP = Inches(6.85)


# ──────────────────────────── 데이터 모델 ────────────────────────────

@dataclass
class CoverSection:
    """표지 슬라이드 데이터."""
    type: Literal["cover"] = "cover"
    subtitle: str = ""  # 부제목 (예: "기술 제안서")


@dataclass
class TocItem:
    """목차 항목."""
    number: str     # "01", "02" ...
    title: str
    page: int = 0


@dataclass
class TocSection:
    """목차 슬라이드 데이터."""
    type: Literal["toc"] = "toc"
    items: list[TocItem] = field(default_factory=list)


@dataclass
class ContentSection:
    """본문 슬라이드 데이터."""
    type: Literal["content"] = "content"
    title: str = ""
    body: list[str] = field(default_factory=list)
    image_path: str | None = None          # 이미지 파일 경로/URL
    image_position: Literal["right", "bottom", "full"] = "right"
    # mermaid 원본 코드. 주어지면 image_path보다 우선해 테마색으로 인라인 렌더링.
    image_mermaid: str | None = None


@dataclass
class ScheduleItem:
    """일정 항목."""
    phase: str        # 단계명
    task: str         # 세부 작업
    duration: str     # "2주", "1개월" 등
    months: list[int] = field(default_factory=list)  # 해당 월 (1~12)


@dataclass
class ScheduleSection:
    """일정 슬라이드 데이터."""
    type: Literal["schedule"] = "schedule"
    title: str = "추진 일정"
    total_months: int = 6
    items: list[ScheduleItem] = field(default_factory=list)


@dataclass
class TeamMember:
    """투입 인력."""
    role: str           # 역할 (PM, PL, 개발자 등)
    name: str
    career_years: int   # 경력 연수
    certification: str = ""  # 자격증
    tasks: str = ""     # 담당 업무


@dataclass
class TeamSection:
    """인력 구성 슬라이드 데이터."""
    type: Literal["team"] = "team"
    title: str = "투입 인력 구성"
    members: list[TeamMember] = field(default_factory=list)


@dataclass
class DataTableSection:
    """범용 데이터 표 슬라이드."""
    type: Literal["data_table"] = "data_table"
    title: str = ""
    table_title: str = ""
    columns: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)


SectionType = CoverSection | TocSection | ContentSection | ScheduleSection | TeamSection | DataTableSection


@dataclass
class ProposalData:
    """제안서 전체 데이터."""
    title: str                  # 사업명
    company: str                # 업체명
    bid_org: str = ""           # 발주기관
    date: str = ""              # 제출일
    sections: list[SectionType] = field(default_factory=list)


# ──────────────────────────── 생성기 클래스 ────────────────────────────

class ProposalPPTGenerator:
    """나라장터 제안서 PPT를 생성한다."""

    def __init__(
        self,
        template_path: str | Path | None = None,
        theme: Theme | str | dict | None = None,
    ):
        if template_path:
            path = Path(template_path)
            if not path.exists():
                raise FileNotFoundError(f"템플릿을 찾을 수 없습니다: {path}")
            self.prs = Presentation(str(path))
        else:
            self.prs = Presentation()
            self.prs.slide_width = SLIDE_WIDTH
            self.prs.slide_height = SLIDE_HEIGHT

        self._blank_layout = self.prs.slide_layouts[6]  # 빈 레이아웃
        self._page_number = 0

        # 테마 주입: Theme 객체 / 이름 문자열 / 인라인 dict 모두 허용
        self.theme = resolve_theme(theme)

    # ──────────── public API ────────────

    def generate(self, proposal_data: ProposalData) -> bytes:
        """제안서 데이터를 받아 PPT 바이트를 반환한다."""
        self._page_number = 0

        for section in proposal_data.sections:
            match section.type:
                case "cover":
                    self._add_cover(proposal_data, section)
                case "toc":
                    self._add_toc(section)
                case "content":
                    self._add_content(section)
                case "schedule":
                    self._add_schedule(section)
                case "team":
                    self._add_team(section)
                case "data_table":
                    self._add_data_table(section)

        buf = io.BytesIO()
        self.prs.save(buf)
        return buf.getvalue()

    # ──────────── 공통 헬퍼 ────────────

    def _new_slide(self, show_page_number: bool = True):
        """빈 슬라이드를 추가하고 반환한다."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._page_number += 1
        if show_page_number and self._page_number > 1:
            self._add_page_number(slide)
        return slide

    def _set_font(
        self,
        run,
        size=None,
        bold: bool = False,
        color=None,
        font_name: str | None = None,
    ):
        """텍스트 런에 테마 폰트/색을 설정한다.

        Latin 폰트(font_name 또는 theme.font_regular/bold)와 별개로
        한글 폴백 폰트(theme.font_korean)를 <a:ea> 요소로 주입해
        영문 디스플레이 폰트가 한글 글리프를 갖지 않더라도 평가위원 PC에서
        한글이 깨지거나 □ 박스로 출력되지 않도록 보장한다.
        """
        latin = font_name or (self.theme.font_bold if bold else self.theme.font_regular)
        run.font.name = latin
        run.font.size = size if size is not None else self.theme.body_size
        run.font.bold = bold
        run.font.color.rgb = color if color is not None else self.theme.color_text
        self._inject_ea_font(run, self.theme.font_korean)

    @staticmethod
    def _inject_ea_font(run, ea_typeface: str):
        """런의 rPr에 <a:ea typeface="..."/> 를 주입(있으면 갱신)한다."""
        if not ea_typeface:
            return
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", ea_typeface)

    @staticmethod
    def _set_cell_border(cell, color: RGBColor, weight_emu: int = 6350):
        """셀 4면 테두리를 단일 색/굵기로 설정한다. weight 6350 EMU = 0.5pt."""
        tcPr = cell._tc.get_or_add_tcPr()
        hex_color = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
        for side in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            existing = tcPr.find(qn(side))
            if existing is not None:
                tcPr.remove(existing)
            ln = etree.SubElement(tcPr, qn(side))
            ln.set("w", str(weight_emu))
            ln.set("cap", "flat")
            ln.set("cmpd", "sng")
            ln.set("algn", "ctr")
            fill = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(fill, qn("a:srgbClr"))
            srgb.set("val", hex_color)

    def _style_table(self, table, *, header_row: int = 0):
        """표 전체에 테마 테두리/셀 padding을 적용한다."""
        rows = len(table.rows)
        cols = len(table.columns)
        for ri in range(rows):
            for ci in range(cols):
                cell = table.cell(ri, ci)
                self._set_cell_border(cell, self.theme.color_divider)
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)

    def _add_textbox(self, slide, left, top, width, height, text: str, **font_kwargs):
        """텍스트박스를 추가하고 런을 반환한다."""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        self._set_font(run, **font_kwargs)
        return txBox, tf, p, run

    def _add_slide_header(self, slide, title: str):
        """슬라이드 상단 제목 + 구분선을 추가한다."""
        self._add_textbox(
            slide, MARGIN_LEFT, HEADER_TOP, CONTENT_WIDTH, Inches(0.6),
            title, size=self.theme.title_size, bold=True, color=self.theme.color_primary,
        )
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_LEFT, Inches(1.15), CONTENT_WIDTH, Pt(2),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.color_primary
        line.line.fill.background()

    def _add_page_number(self, slide):
        """우측 하단 페이지 번호."""
        _, _, _, run = self._add_textbox(
            slide,
            Inches(9.0), FOOTER_TOP, Inches(1.0), Inches(0.3),
            str(self._page_number - 1),  # 표지 제외
            size=self.theme.page_num_size, color=self.theme.color_text_muted,
        )
        run.font.name = self.theme.font_regular

    def _add_footer_line(self, slide):
        """하단 구분선."""
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            MARGIN_LEFT, Inches(6.75), CONTENT_WIDTH, Pt(1),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.color_divider
        line.line.fill.background()

    # ──────────── 표지 (Cover) ────────────

    def _add_cover(self, data: ProposalData, section: CoverSection):
        if self.theme.cover_layout == "asymmetric":
            return self._add_cover_asymmetric(data, section)
        return self._add_cover_accent_bar(data, section)

    def _add_cover_accent_bar(self, data: ProposalData, section: CoverSection):
        slide = self._new_slide(show_page_number=False)

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.theme.color_bg_accent

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.8), Inches(0.08), Inches(2.8),
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme.color_bar_on_accent
        bar.line.fill.background()

        if section.subtitle:
            self._add_textbox(
                slide, Inches(1.0), Inches(1.9), Inches(8.0), Inches(0.5),
                section.subtitle,
                size=self.theme.cover_subtitle_size, color=self.theme.color_text_on_accent,
            )

        self._add_textbox(
            slide, Inches(1.0), Inches(2.5), Inches(8.0), Inches(1.5),
            data.title,
            size=self.theme.cover_title_size, bold=True, color=self.theme.color_text_on_accent,
        )

        meta_lines = []
        if data.bid_org:
            meta_lines.append(f"발주기관  |  {data.bid_org}")
        if data.date:
            meta_lines.append(f"제 출 일  |  {data.date}")
        if meta_lines:
            self._add_textbox(
                slide, Inches(1.0), Inches(4.3), Inches(8.0), Inches(0.8),
                "\n".join(meta_lines),
                size=self.theme.cover_meta_size, color=self.theme.color_text_on_accent,
            )

        _, tf, p, run = self._add_textbox(
            slide, Inches(1.0), Inches(5.8), Inches(8.0), Inches(0.5),
            data.company,
            size=self.theme.cover_company_size, bold=True, color=self.theme.color_text_on_accent,
        )
        p.alignment = PP_ALIGN.RIGHT

    def _add_cover_asymmetric(self, data: ProposalData, section: CoverSection):
        """좌상단 거대 타이포 + 우하단 모노 메타. 다크 풀블리드 배경."""
        slide = self._new_slide(show_page_number=False)

        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.theme.color_bg_accent

        # 상단 모노 마크
        self._add_textbox(
            slide, Inches(0.6), Inches(0.45), Inches(8), Inches(0.35),
            "PROPOSAL  //  " + (data.bid_org or "PUBLIC BID"),
            size=Pt(9), color=self.theme.color_text_on_accent,
            font_name=self.theme.font_display,
        )

        # 좌상단 부제 (작게)
        if section.subtitle:
            self._add_textbox(
                slide, Inches(0.6), Inches(1.0), Inches(8), Inches(0.4),
                section.subtitle,
                size=self.theme.cover_subtitle_size, color=self.theme.color_text_on_accent,
            )

        # 거대 타이틀 (좌측 정렬, 다중행 허용)
        self._add_textbox(
            slide, Inches(0.6), Inches(1.7), Inches(11.6), Inches(3.3),
            data.title,
            size=self.theme.cover_title_size, bold=True, color=self.theme.color_text_on_accent,
        )

        # 좌하단: 회사명
        self._add_textbox(
            slide, Inches(0.6), Inches(5.6), Inches(7), Inches(0.6),
            data.company,
            size=self.theme.cover_company_size, bold=True, color=self.theme.color_text_on_accent,
        )

        # 우하단: 메타 (모노, 우측정렬)
        meta_lines = []
        if data.date:
            meta_lines.append(f"DATE  //  {data.date}")
        if data.bid_org:
            meta_lines.append(f"CLIENT  //  {data.bid_org}")
        if meta_lines:
            _, tf, p, _ = self._add_textbox(
                slide, Inches(7.0), Inches(5.5), Inches(5.2), Inches(1.0),
                "\n".join(meta_lines),
                size=self.theme.cover_meta_size, color=self.theme.color_text_on_accent,
                font_name=self.theme.font_display,
            )
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.RIGHT

        # 하단 풀폭 구분선
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(6.6), Inches(11.6), Pt(1),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.theme.color_bar_on_accent
        line.line.fill.background()

    # ──────────── 목차 (TOC) ────────────

    def _add_toc(self, section: TocSection):
        slide = self._new_slide()
        self._add_slide_header(slide, "목 차")

        y = Inches(1.7)
        row_height = Inches(0.65)

        for item in section.items:
            # 번호 원형 배지
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(1.0), y + Inches(0.08), Inches(0.42), Inches(0.42),
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = self.theme.color_primary
            circle.line.fill.background()
            tf = circle.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            run = tf.paragraphs[0].add_run()
            run.text = item.number
            self._set_font(run, size=Pt(11), bold=True, color=self.theme.color_text_on_accent)

            self._add_textbox(
                slide, Inches(1.65), y, Inches(6.5), Inches(0.5),
                item.title,
                size=Pt(14), bold=True, color=self.theme.color_text,
            )

            if item.page:
                _, _, p, _ = self._add_textbox(
                    slide, Inches(8.5), y, Inches(1.0), Inches(0.5),
                    str(item.page),
                    size=Pt(12), color=self.theme.color_text_muted,
                )
                p.alignment = PP_ALIGN.RIGHT

            if item != section.items[-1]:
                dotline = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1.0), y + Inches(0.58), Inches(8.5), Pt(0.5),
                )
                dotline.fill.solid()
                dotline.fill.fore_color.rgb = self.theme.color_divider
                dotline.line.fill.background()

            y += row_height

        self._add_footer_line(slide)

    # ──────────── 이미지 다운로드 ────────────

    @staticmethod
    def _resolve_image(image_path: str | None) -> str | None:
        """이미지 경로를 실제 사용 가능한 로컬 경로로 반환한다.

        - 로컬 파일이면 그대로 반환
        - URL이면 다운로드하여 임시 파일 경로 반환
        - 실패하면 None
        """
        if not image_path:
            return None

        # 로컬 파일
        if not image_path.startswith(("http://", "https://")):
            return image_path if Path(image_path).exists() else None

        # URL → 다운로드
        try:
            resp = httpx.get(image_path, timeout=30, follow_redirects=True)
            resp.raise_for_status()

            ct = resp.headers.get("content-type", "")
            ext = ".png" if "png" in ct else ".jpg"
            tmp = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
            tmp.write(resp.content)
            tmp.close()
            return tmp.name
        except Exception as e:
            print(f"[PPT] 이미지 다운로드 실패: {e}")
            return None

    # ──────────── 마크다운 파싱 ────────────

    def _add_rich_text(self, paragraph, text: str, base_size=None, base_color=None):
        """**bold** 마크다운을 파싱하여 런을 분리 추가한다."""
        size = base_size if base_size is not None else self.theme.body_size
        color = base_color if base_color is not None else self.theme.color_text
        parts = re.split(r'(\*\*.*?\*\*)', text)
        for part in parts:
            if not part:
                continue
            run = paragraph.add_run()
            if part.startswith('**') and part.endswith('**'):
                run.text = part[2:-2]
                self._set_font(run, size=size, bold=True, color=color)
            else:
                run.text = part
                self._set_font(run, size=size, color=color)

    # ──────────── 본문 (Content) ────────────

    # 레이아웃별 최대 시각적 본문 줄 수 (줄바꿈 후 기준, 잘림 방지)
    _MAX_VISUAL_LINES_TEXT_ONLY = 12
    _MAX_VISUAL_LINES_IMAGE_RIGHT = 11
    _MAX_VISUAL_LINES_IMAGE_BOTTOM = 5

    # 레이아웃별 한 줄에 들어가는 한국어 글자 수 추정치
    _CHARS_PER_LINE_TEXT_ONLY = 40
    _CHARS_PER_LINE_IMAGE_RIGHT = 22  # 좌측 5인치 폭
    _CHARS_PER_LINE_IMAGE_BOTTOM = 40

    @staticmethod
    def _estimate_visual_lines(line: str, chars_per_line: int) -> int:
        """렌더링 시 차지할 시각적 줄 수를 추정한다.

        - 빈 줄: 1줄
        - 하위 불릿("  • "): 들여쓰기로 가용 폭이 좁으므로 chars_per_line의 85%
        - 일반: 글자수 / chars_per_line 올림
        - 굵게(**) 표기는 글자수에서 제외
        """
        if not line:
            return 1
        # ** 마커는 길이에서 제외
        text = re.sub(r'\*\*', '', line)
        # 하위 불릿은 폭이 좁음
        is_sub = line.startswith("  • ")
        cpl = max(10, int(chars_per_line * 0.85)) if is_sub else chars_per_line
        return max(1, -(-len(text) // cpl))  # 올림 나눗셈

    @classmethod
    def _pick_body_style(
        cls,
        chunk: list[str],
        chars_per_line: int,
    ) -> tuple[Pt, Pt, Pt, Pt]:
        """청크의 콘텐츠 밀도에 따라 (본문, 소제목, line_spacing, space_after)를 결정한다.

        밀도가 높을수록 폰트를 줄여 잘림을 방지한다 (14pt → 12pt → 10pt 단계).
        """
        total = sum(cls._estimate_visual_lines(l, chars_per_line) for l in chunk)
        if total <= 6:
            return Pt(14), Pt(15), Pt(22), Pt(10)  # 여유: 키움
        if total <= 10:
            return Pt(12), Pt(13), Pt(20), Pt(8)   # 기본
        return Pt(10), Pt(12), Pt(16), Pt(6)        # 빽빽: 줄임

    @classmethod
    def _chunk_body(
        cls,
        body: list[str],
        max_visual_lines: int,
        chars_per_line: int,
    ) -> list[list[str]]:
        """본문 줄 목록을 시각적 줄 수 기준으로 분할한다.

        - 한 항목의 줄바꿈 후 줄 수까지 누적해 max_visual_lines를 초과하면 분할.
        - 하위 불릿("  • ")이 부모("■")와 분리되지 않도록 분할 지점을 보정한다.
        """
        if not body or max_visual_lines <= 0:
            return [body]

        def _is_heading(s: str) -> bool:
            return s.startswith("■ ") or s.startswith("● ")

        chunks: list[list[str]] = []
        current: list[str] = []
        used = 0
        for line in body:
            cost = cls._estimate_visual_lines(line, chars_per_line)
            # 현재 청크가 비어있지 않은데 추가하면 한도 초과 → 새 청크 시작
            if current and used + cost > max_visual_lines:
                # 다음 줄이 하위 불릿이면 직전 부모(■/●/일반)도 함께 다음 청크로 이동
                if line.startswith("  • "):
                    detach_from = len(current)
                    while detach_from > 0 and current[detach_from - 1].startswith("  • "):
                        detach_from -= 1
                    if detach_from > 0:
                        detach_from -= 1  # 부모 라인 자체 포함
                    if detach_from > 0:
                        carry = current[detach_from:]
                        current = current[:detach_from]
                        chunks.append(current)
                        current = carry + [line]
                        used = sum(cls._estimate_visual_lines(l, chars_per_line) for l in current)
                        continue
                # 현재 청크가 헤딩(■/●)으로 끝나면 그 헤딩을 다음 청크로 이동
                # (헤딩만 남은 빈 슬라이드 방지)
                trailing_headings = 0
                i = len(current) - 1
                while i >= 0 and _is_heading(current[i]):
                    trailing_headings += 1
                    i -= 1
                if trailing_headings > 0 and trailing_headings < len(current):
                    carry = current[-trailing_headings:]
                    current = current[:-trailing_headings]
                    chunks.append(current)
                    current = carry + [line]
                    used = sum(cls._estimate_visual_lines(l, chars_per_line) for l in current)
                    continue
                chunks.append(current)
                current = [line]
                used = cost
            else:
                current.append(line)
                used += cost
        # 마지막 청크가 헤딩만 들어있으면 직전 청크에 합치기
        if current:
            if all(_is_heading(l) for l in current) and chunks:
                chunks[-1].extend(current)
            else:
                chunks.append(current)
        return chunks if chunks else [body]

    def _theme_mermaid_colors(self) -> dict:
        """현재 테마에서 mermaid themeVariables용 색 팔레트를 추출한다."""
        t = self.theme
        def hexstr(c) -> str:
            return f"#{c[0]:02X}{c[1]:02X}{c[2]:02X}"
        return {
            "primary": hexstr(t.color_primary),
            "text": hexstr(t.color_text),
            "background": hexstr(t.color_bg_body),
            "line": hexstr(t.color_primary),
            "secondary": hexstr(t.color_light_surface),
        }

    def _render_mermaid_to_tempfile(self, code: str) -> str | None:
        """mermaid 코드를 테마색으로 PNG 렌더링 후 임시파일 경로 반환. 실패 시 None."""
        try:
            from .mermaid_renderer import render_mermaid_to_png
            png = render_mermaid_to_png(
                code,
                width=1600,
                height=900,
                background="white",
                theme="neutral",
                theme_colors=self._theme_mermaid_colors(),
            )
            tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            tmp.write(png)
            tmp.close()
            return tmp.name
        except Exception as e:
            print(f"[PPT] 인라인 mermaid 렌더 실패: {e}")
            return None

    def _add_content(self, section: ContentSection):
        # mermaid 인라인 코드가 있으면 테마색으로 즉시 렌더 (image_path보다 우선)
        if section.image_mermaid:
            rendered = self._render_mermaid_to_tempfile(section.image_mermaid)
            if rendered:
                resolved_image = rendered
            else:
                resolved_image = self._resolve_image(section.image_path)
        else:
            resolved_image = self._resolve_image(section.image_path)

        # 전체 이미지 슬라이드는 텍스트가 없으므로 분할 불필요
        if resolved_image and section.image_position == "full":
            slide = self._new_slide()
            self._add_slide_header(slide, section.title)
            self._add_content_image_full(slide, resolved_image)
            self._add_footer_line(slide)
            return

        # 레이아웃별 분할 한도 결정
        if resolved_image and section.image_position == "right":
            max_lines = self._MAX_VISUAL_LINES_IMAGE_RIGHT
            chars_per_line = self._CHARS_PER_LINE_IMAGE_RIGHT
        elif resolved_image and section.image_position == "bottom":
            max_lines = self._MAX_VISUAL_LINES_IMAGE_BOTTOM
            chars_per_line = self._CHARS_PER_LINE_IMAGE_BOTTOM
        else:
            max_lines = self._MAX_VISUAL_LINES_TEXT_ONLY
            chars_per_line = self._CHARS_PER_LINE_TEXT_ONLY

        chunks = self._chunk_body(section.body, max_lines, chars_per_line)

        for idx, chunk in enumerate(chunks):
            slide = self._new_slide()
            title = section.title if idx == 0 else f"{section.title} (계속)"
            self._add_slide_header(slide, title)

            sub_section = ContentSection(
                title=title,
                body=chunk,
                image_path=section.image_path,
                image_position=section.image_position,
            )

            body_size, heading_size, line_sp, space_after = self._pick_body_style(
                chunk, chars_per_line
            )

            if resolved_image and section.image_position == "right":
                self._add_content_with_image_right(
                    slide, sub_section, resolved_image,
                    body_size=body_size, line_sp=line_sp, space_after=space_after,
                )
            elif resolved_image and section.image_position == "bottom":
                self._add_content_with_image_bottom(
                    slide, sub_section, resolved_image,
                    body_size=body_size, line_sp=line_sp, space_after=space_after,
                )
            else:
                self._add_content_text_only(
                    slide, sub_section,
                    body_size=body_size, heading_size=heading_size,
                    line_sp=line_sp, space_after=space_after,
                )

            self._add_footer_line(slide)

    def _add_content_text_only(
        self, slide, section: ContentSection,
        body_size: Pt = BODY_SIZE, heading_size: Pt = Pt(13),
        line_sp: Pt = Pt(20), space_after: Pt = Pt(8),
    ):
        """텍스트만 있는 본문 슬라이드."""
        txBox = slide.shapes.add_textbox(
            MARGIN_LEFT, BODY_TOP, CONTENT_WIDTH, BODY_HEIGHT,
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(section.body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = space_after
            p.line_spacing = line_sp

            if not line:
                continue

            if line.startswith("■"):
                self._add_rich_text(p, line, base_size=heading_size, base_color=self.theme.color_primary)
            elif line.startswith("●"):
                self._add_rich_text(p, line, base_size=body_size, base_color=self.theme.color_accent)
            elif line.startswith("  • "):
                p.level = 1
                self._add_rich_text(p, line, base_size=body_size, base_color=self.theme.color_text)
            else:
                self._add_rich_text(p, f"  {line}", base_size=body_size, base_color=self.theme.color_text)

    def _add_content_with_image_right(
        self, slide, section: ContentSection, image_path: str,
        body_size: Pt = BODY_SIZE, line_sp: Pt = Pt(20), space_after: Pt = Pt(8),
    ):
        """좌측 텍스트 + 우측 이미지."""
        # 텍스트 (좌측 60%)
        txBox = slide.shapes.add_textbox(
            MARGIN_LEFT, BODY_TOP, Inches(5.0), BODY_HEIGHT,
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(section.body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = space_after
            p.line_spacing = line_sp
            if line:
                self._add_rich_text(p, f"• {line}", base_size=body_size, base_color=self.theme.color_text)

        # 이미지 (우측 40%)
        slide.shapes.add_picture(
            image_path,
            Inches(6.0), BODY_TOP, Inches(3.6), Inches(3.5),
        )

    def _add_content_with_image_bottom(
        self, slide, section: ContentSection, image_path: str,
        body_size: Pt = BODY_SIZE, line_sp: Pt = Pt(20), space_after: Pt = Pt(6),
    ):
        """상단 텍스트 + 하단 이미지."""
        txBox = slide.shapes.add_textbox(
            MARGIN_LEFT, BODY_TOP, CONTENT_WIDTH, Inches(2.2),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(section.body):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = space_after
            p.line_spacing = line_sp
            if line:
                self._add_rich_text(p, f"• {line}", base_size=body_size, base_color=self.theme.color_text)

        slide.shapes.add_picture(
            image_path,
            Inches(1.5), Inches(4.0), Inches(7.3), Inches(2.5),
        )

    def _add_content_image_full(self, slide, image_path: str):
        """전체 이미지 슬라이드."""
        slide.shapes.add_picture(
            image_path,
            MARGIN_LEFT, BODY_TOP, CONTENT_WIDTH, BODY_HEIGHT,
        )

    # ──────────── 추진 일정 (Schedule) ────────────

    def _add_schedule(self, section: ScheduleSection):
        slide = self._new_slide()
        self._add_slide_header(slide, section.title)

        total = section.total_months
        table_left = MARGIN_LEFT
        table_top = Inches(1.7)
        col_widths = [Inches(1.5), Inches(2.5)] + [
            Inches(4.9 / total) for _ in range(total)
        ]
        rows = len(section.items) + 1  # 헤더 + 데이터

        table_shape = slide.shapes.add_table(
            rows, 2 + total,
            table_left, table_top,
            sum(col_widths, Emu(0)), Inches(0.4 * rows),
        )
        table = table_shape.table
        self._style_table(table)

        # 열 너비 설정
        for i, w in enumerate(col_widths):
            table.columns[i].width = int(w)

        # 헤더 행
        headers = ["단계", "세부 작업"] + [f"{m}월" for m in range(1, total + 1)]
        for ci, text in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.color_table_header_bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run, size=Pt(10), bold=True, color=self.theme.color_table_header_text)

        # 데이터 행
        for ri, item in enumerate(section.items, start=1):
            # 단계
            c_phase = table.cell(ri, 0)
            c_phase.text = ""
            c_phase.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c_phase.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item.phase
            self._set_font(run, size=Pt(10), bold=True, color=self.theme.color_text)

            # 세부 작업
            c_task = table.cell(ri, 1)
            c_task.text = ""
            c_task.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c_task.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = item.task
            self._set_font(run, size=Pt(10), color=self.theme.color_text)

            # 간트 바 (해당 월 셀에 색칠)
            for m in item.months:
                if 1 <= m <= total:
                    cell = table.cell(ri, 1 + m)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme.color_accent

            # 짝수 행 배경 (간트 바 셀은 제외)
            if ri % 2 == 0:
                gantt_cols = {1 + m for m in item.months if 1 <= m <= total}
                for ci in range(2 + total):
                    if ci not in gantt_cols:
                        cell = table.cell(ri, ci)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.theme.color_light_surface

        self._add_footer_line(slide)

    # ──────────── 투입 인력 (Team) ────────────

    def _add_team(self, section: TeamSection):
        slide = self._new_slide()
        self._add_slide_header(slide, section.title)

        cols = 5
        rows = len(section.members) + 1
        col_widths = [Inches(1.2), Inches(1.2), Inches(1.0), Inches(2.0), Inches(3.5)]

        table_shape = slide.shapes.add_table(
            rows, cols,
            MARGIN_LEFT, Inches(1.7),
            sum(col_widths, Emu(0)), Inches(0.45 * rows),
        )
        table = table_shape.table
        self._style_table(table)

        for i, w in enumerate(col_widths):
            table.columns[i].width = int(w)

        # 헤더
        headers = ["역할", "성명", "경력", "자격/학위", "담당 업무"]
        for ci, text in enumerate(headers):
            cell = table.cell(0, ci)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.color_table_header_bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run, size=Pt(10), bold=True, color=self.theme.color_table_header_text)

        # 데이터
        for ri, member in enumerate(section.members, start=1):
            values = [
                member.role,
                member.name,
                f"{member.career_years}년",
                member.certification,
                member.tasks,
            ]
            for ci, val in enumerate(values):
                cell = table.cell(ri, ci)
                cell.text = ""
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if ci < 3 else PP_ALIGN.LEFT
                run = p.add_run()
                run.text = val
                self._set_font(run, size=Pt(10), color=self.theme.color_text)

                # 역할 열 강조
                if ci == 0:
                    self._set_font(run, size=Pt(10), bold=True, color=self.theme.color_primary)

            # 짝수 행 배경
            if ri % 2 == 0:
                for ci in range(cols):
                    cell = table.cell(ri, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme.color_light_surface

        self._add_footer_line(slide)

    # ──────────── 범용 데이터 표 (DataTable) ────────────

    def _add_data_table(self, section: DataTableSection):
        slide = self._new_slide()
        self._add_slide_header(slide, section.title)

        # 표 소제목
        table_top = Inches(1.5)
        if section.table_title:
            self._add_textbox(
                slide, MARGIN_LEFT, Inches(1.3), CONTENT_WIDTH, Inches(0.4),
                section.table_title, size=Pt(13), bold=True, color=self.theme.color_primary,
            )
            table_top = Inches(1.8)

        cols = len(section.columns)
        if cols == 0:
            return

        data_rows = section.rows
        rows = len(data_rows) + 1  # 헤더 + 데이터

        # 행이 너무 많으면 제한
        max_rows = 12
        if rows > max_rows + 1:
            data_rows = data_rows[:max_rows]
            rows = max_rows + 1

        col_width = int(CONTENT_WIDTH / cols)
        row_height = min(Inches(0.4), int(Inches(4.5) / rows))

        table_shape = slide.shapes.add_table(
            rows, cols,
            MARGIN_LEFT, table_top,
            CONTENT_WIDTH, row_height * rows,
        )
        table = table_shape.table
        self._style_table(table)

        for ci in range(cols):
            table.columns[ci].width = col_width

        # 헤더
        for ci, text in enumerate(section.columns):
            cell = table.cell(0, ci)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme.color_table_header_bg
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            self._set_font(run, size=Pt(10), bold=True, color=self.theme.color_table_header_text)

        # 데이터
        for ri, row_data in enumerate(data_rows, start=1):
            for ci in range(cols):
                cell_text = row_data[ci] if ci < len(row_data) else ""
                cell = table.cell(ri, ci)
                cell.text = ""
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER if ci == 0 else PP_ALIGN.LEFT
                run = p.add_run()
                run.text = cell_text
                self._set_font(run, size=Pt(9), color=self.theme.color_text)

                if ci == 0:
                    self._set_font(run, size=Pt(9), bold=True, color=self.theme.color_primary)

            if ri % 2 == 0:
                for ci in range(cols):
                    cell = table.cell(ri, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme.color_light_surface

        self._add_footer_line(slide)
